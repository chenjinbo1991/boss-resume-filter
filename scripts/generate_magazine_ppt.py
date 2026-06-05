"""
BOSS 招聘系统操作指南 — 杂志风格 PPT 生成脚本
用法: python scripts/generate_magazine_ppt.py
输出: docs/BOSS招聘系统操作指南.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── 设计系统 ───────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色
C_DARK    = RGBColor(0x1B, 0x28, 0x38)  # 深蓝墨
C_ACCENT  = RGBColor(0xE8, 0x65, 0x2C)  # 亮橙
C_BLUE    = RGBColor(0x2D, 0x7D, 0xD2)  # 中蓝
C_LIGHT   = RGBColor(0xF0, 0xF2, 0xF5)  # 浅灰底
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0x6B, 0x7B, 0x8D)  # 正文灰
C_DGRAY   = RGBColor(0x3D, 0x4F, 0x5F)  # 深灰
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)
C_RED     = RGBColor(0xE7, 0x4C, 0x3C)
C_YELLOW  = RGBColor(0xF3, 0x9C, 0x12)

# 字体
FONT_TITLE  = "Microsoft YaHei"
FONT_BODY   = "Microsoft YaHei"

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS_DIR = os.path.join(BASE_DIR, "docs", "assets", "user-guide")
OUTPUT_PATH = os.path.join(BASE_DIR, "docs", "BOSS招聘系统操作指南.pptx")


# ─── 工具函数 ───────────────────────────────────────────────
def set_font(run, size=18, bold=False, color=C_DARK, name=FONT_BODY, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    run.font.italic = italic


def add_textbox(slide, left, top, width, height, text="", size=18, bold=False,
                color=C_DARK, alignment=PP_ALIGN.LEFT, name=FONT_BODY, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color, name, italic)
    return txBox, tf


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """安全添加图片，文件不存在时返回 None"""
    if os.path.exists(img_path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(img_path, **kwargs)
    return None


def add_page_number(slide, num, total=15):
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
                f"{num} / {total}", size=10, color=C_GRAY, alignment=PP_ALIGN.RIGHT)


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6)):
    add_rect(slide, left, top, width, height, C_ACCENT)


def add_section_tag(slide, text, left, top):
    """添加小号分类标签"""
    shape = add_rounded_rect(slide, left, top, Inches(1.6), Inches(0.35), C_ACCENT)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, 10, bold=True, color=C_WHITE)
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def add_card(slide, left, top, width, height, fill=C_WHITE, shadow=True):
    """添加卡片背景"""
    shape = add_rounded_rect(slide, left, top, width, height, fill)
    if shadow:
        shape.shadow.inherit = False
    return shape


# ─── 页面背景 ───────────────────────────────────────────────
def set_bg_white(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_WHITE


def set_bg_dark(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK


def set_bg_light(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_LIGHT


# ─── 幻灯片生成 ──────────────────────────────────────────────

def slide_01_cover(prs):
    """封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg_dark(slide)

    # 左侧装饰色块
    add_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, C_ACCENT)

    # 右上角装饰圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x24, 0x34, 0x47)
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(5.0), Inches(3), Inches(3))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x24, 0x34, 0x47)
    circle2.line.fill.background()

    # 分类标签
    add_section_tag(slide, "操作指南 v2.9", Inches(1.2), Inches(1.5))

    # 主标题
    add_textbox(slide, Inches(1.2), Inches(2.2), Inches(9), Inches(1.5),
                "BOSS 招聘系统", size=54, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(1.2), Inches(3.5), Inches(9), Inches(1.2),
                "操作指南", size=48, bold=True, color=C_ACCENT)

    # 副标题
    add_textbox(slide, Inches(1.2), Inches(5.0), Inches(8), Inches(0.5),
                "岗位配置 · 智能筛选 · AI 评估 · 自动打招呼 · 结果导出",
                size=16, color=C_GRAY)

    # 底部信息
    add_textbox(slide, Inches(1.2), Inches(6.5), Inches(4), Inches(0.4),
                "适用对象：图形界面用户", size=12, color=C_GRAY)
    add_textbox(slide, Inches(9), Inches(6.5), Inches(3.5), Inches(0.4),
                "2026 · 内部培训资料", size=12, color=C_GRAY, alignment=PP_ALIGN.RIGHT)


def slide_02_overview(prs):
    """系统能做什么"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(slide)

    # 左侧色带
    add_rect(slide, Inches(0), Inches(0), Inches(5.0), SLIDE_H, C_DARK)

    # 左半标题
    add_section_tag(slide, "OVERVIEW", Inches(0.8), Inches(1.2))
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(3.5), Inches(1.5),
                "系统能\n做什么", size=44, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(3.5), Inches(1.5),
                "把候选人获取、筛选、评分、打招呼和导出串成一个可重复执行的流程",
                size=14, color=C_GRAY, italic=False)

    # 右半 — 6 个核心能力卡片
    capabilities = [
        ("01", "规则筛选", "学历 / 经验 / 年龄 / 薪资 / 地点 / 必要条件"),
        ("02", "技能评分", "关键词权重匹配，四维模型打分"),
        ("03", "AI 二次评估", "大模型对通过筛选者做 ±10 分调整"),
        ("04", "自动滚动提取", "智能滚动 + 批量提取候选人卡片"),
        ("05", "自动打招呼", "按推荐等级阈值自动发送沟通消息"),
        ("06", "导出统计", "Excel 导出 + 按岗位统计看板"),
    ]
    start_y = 0.7
    for i, (num, title, desc) in enumerate(capabilities):
        y = start_y + i * 1.05
        # 编号圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(y), Inches(0.5), Inches(0.5))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_ACCENT
        dot.line.fill.background()
        tf = dot.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, 11, bold=True, color=C_WHITE)

        add_textbox(slide, Inches(6.3), Inches(y - 0.05), Inches(3), Inches(0.35),
                    title, size=16, bold=True, color=C_DARK)
        add_textbox(slide, Inches(6.3), Inches(y + 0.3), Inches(6.5), Inches(0.5),
                    desc, size=11, color=C_GRAY)

    add_page_number(slide, 2)


def slide_03_workflow(prs):
    """工作流程全景"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_light(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "工作流程全景", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(8), Inches(0.4),
                "从配置到导出，7 步完成完整筛选流程", size=14, color=C_GRAY)

    steps = [
        ("1", "启动系统", "双击打开程序", C_DARK),
        ("2", "配置岗位", "新建岗位规则\n或解析 JD", C_BLUE),
        ("3", "配置 AI", "可选\n配置大模型", RGBColor(0x6C, 0x5C, 0xE7)),
        ("4", "连接浏览器", "连接 Chrome\n打开推荐页", C_ACCENT),
        ("5", "运行筛选", "滚动提取 +\n评分 + AI", C_GREEN),
        ("6", "查看结果", "候选人列表\n右键操作", RGBColor(0x00, 0xB8, 0x94)),
        ("7", "导出/统计", "Excel 导出\n数据统计", C_DARK),
    ]

    start_x = 0.5
    card_w = 1.6
    gap = 0.22
    y_pos = 2.2

    for i, (num, title, desc, color) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        # 卡片
        card = add_rounded_rect(slide, Inches(x), Inches(y_pos), Inches(card_w), Inches(3.5), C_WHITE)
        # 顶部色条
        add_rect(slide, Inches(x), Inches(y_pos), Inches(card_w), Inches(0.08), color)
        # 编号
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                            Inches(x + 0.55), Inches(y_pos + 0.3), Inches(0.5), Inches(0.5))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = color
        num_circle.line.fill.background()
        tf = num_circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, 16, bold=True, color=C_WHITE)

        add_textbox(slide, Inches(x + 0.1), Inches(y_pos + 1.0), Inches(card_w - 0.2), Inches(0.4),
                    title, size=14, bold=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.1), Inches(y_pos + 1.5), Inches(card_w - 0.2), Inches(1.5),
                    desc, size=10, color=C_GRAY, alignment=PP_ALIGN.CENTER)

        # 箭头连接线（除了最后一个）
        if i < len(steps) - 1:
            arrow_x = x + card_w + 0.02
            add_textbox(slide, Inches(arrow_x), Inches(y_pos + 1.3), Inches(0.2), Inches(0.4),
                        "→", size=18, bold=True, color=C_ACCENT, alignment=PP_ALIGN.CENTER)

    # 底部截图缩略
    img_path = os.path.join(ASSETS_DIR, "01-home.png")
    add_image_safe(slide, img_path, Inches(1.5), Inches(6.0), width=Inches(3.0))

    add_page_number(slide, 3)


def slide_04_navigation(prs):
    """主界面导航"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "主界面导航", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(8), Inches(0.4),
                "左侧固定导航栏 + 右侧内容区，6 个功能模块覆盖全流程", size=14, color=C_GRAY)

    # 左侧截图
    img_path = os.path.join(ASSETS_DIR, "01-home.png")
    add_image_safe(slide, img_path, Inches(0.8), Inches(1.8), width=Inches(5.5))

    # 右侧导航卡片
    nav_items = [
        ("🏠", "首页", "全局统计 + 快捷入口", C_DARK),
        ("⚙", "岗位配置", "新建/修改/导入/导出规则", C_BLUE),
        ("▶", "运行控制", "连接浏览器 + 运行参数", C_ACCENT),
        ("📋", "筛选结果", "候选人列表 + 右键操作", C_GREEN),
        ("📊", "数据统计", "按岗位/时间统计", RGBColor(0x6C, 0x5C, 0xE7)),
        ("🔧", "系统设置", "AI 服务商 + API Key", C_GRAY),
    ]

    for i, (icon, title, desc, color) in enumerate(nav_items):
        y = 1.9 + i * 0.85
        # 色块标识
        indicator = add_rect(slide, Inches(6.8), Inches(y), Inches(0.06), Inches(0.6), color)
        add_textbox(slide, Inches(7.1), Inches(y - 0.05), Inches(2.5), Inches(0.35),
                    title, size=16, bold=True, color=C_DARK)
        add_textbox(slide, Inches(7.1), Inches(y + 0.3), Inches(5.5), Inches(0.35),
                    desc, size=12, color=C_GRAY)

    # 底部提示
    add_rounded_rect(slide, Inches(6.8), Inches(6.8), Inches(5.8), Inches(0.5), C_LIGHT)
    add_textbox(slide, Inches(7.0), Inches(6.85), Inches(5.5), Inches(0.4),
                "💡 使用顺序：岗位配置 → 系统设置 → 运行控制 → 筛选结果 → 数据统计",
                size=11, color=C_DGRAY)

    add_page_number(slide, 4)


def slide_05_job_config(prs):
    """岗位配置"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "岗位配置", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.4),
                "岗位规则直接决定候选人是否被淘汰，尤其是必要条件和工作地点", size=14, color=C_GRAY)

    # 左半 — 截图
    img_path = os.path.join(ASSETS_DIR, "02-job-config-full.png")
    add_image_safe(slide, img_path, Inches(0.6), Inches(1.8), width=Inches(6.5))

    # 右半 — 字段说明
    fields = [
        ("岗位名称", "当前岗位规则名称", C_DARK),
        ("最低学历", "低于该学历直接淘汰", C_BLUE),
        ("最低经验", "工作年限不足直接淘汰", C_BLUE),
        ("最大年龄", "超过上限直接淘汰，留空不限", C_ACCENT),
        ("薪资范围", "期望薪资明显超出预算时淘汰", C_YELLOW),
        ("工作地点", "支持多地点: 南京、上海/杭州", C_GREEN),
        ("技能关键词", "用于评分，不一定直接淘汰", RGBColor(0x6C, 0x5C, 0xE7)),
        ("必要条件", "不满足直接淘汰，如\"统招本科\"", C_RED),
    ]

    add_textbox(slide, Inches(7.5), Inches(1.7), Inches(5), Inches(0.4),
                "核心字段", size=18, bold=True, color=C_DARK)

    for i, (name, desc, color) in enumerate(fields):
        y = 2.3 + i * 0.58
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(y + 0.05), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_textbox(slide, Inches(7.8), Inches(y - 0.05), Inches(2.0), Inches(0.3),
                    name, size=12, bold=True, color=C_DARK)
        add_textbox(slide, Inches(9.6), Inches(y - 0.05), Inches(3.2), Inches(0.4),
                    desc, size=11, color=C_GRAY)

    # 权重建议
    add_rounded_rect(slide, Inches(7.5), Inches(7.0), Inches(5.2), Inches(0.35), C_LIGHT)
    add_textbox(slide, Inches(7.7), Inches(7.0), Inches(5.0), Inches(0.35),
                "权重建议：1=普通  2=核心  3=强核心（不要全部设高权重）",
                size=10, color=C_DGRAY)

    add_page_number(slide, 5)


def slide_06_jd_parse(prs):
    """JD 解析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_light(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "招聘需求自动解析", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.4),
                "粘贴 JD → 系统自动提取字段 → 人工检查修正 → 保存", size=14, color=C_GRAY)

    # 四个步骤卡片
    steps = [
        ("1", "粘贴招聘需求", "把 JD 文本粘贴到\n文本框中"),
        ("2", "参考示例模板", "不清楚格式？点击\n\"招聘需求示例\""),
        ("3", "解析 + 检查", "检查自动提取的字段\n手工修正偏差"),
        ("4", "保存配置", "确认无误后保存\n生成岗位规则"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = 0.8 + i * 3.1
        card = add_card(slide, Inches(x), Inches(2.0), Inches(2.8), Inches(2.8))
        # 顶部编号
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                           Inches(x + 1.05), Inches(2.3), Inches(0.7), Inches(0.7))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = C_ACCENT
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, 20, bold=True, color=C_WHITE)

        add_textbox(slide, Inches(x + 0.2), Inches(3.2), Inches(2.4), Inches(0.4),
                    title, size=15, bold=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(3.7), Inches(2.4), Inches(0.9),
                    desc, size=11, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # 底部警告
    warn_shape = add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.7), C_WHITE)
    add_textbox(slide, Inches(1.2), Inches(5.35), Inches(11.0), Inches(0.6),
                "⚠️  解析功能是辅助，不是最终裁决。保存前必须人工检查：薪资上下限、必要条件、技能关键词权重。",
                size=13, bold=True, color=C_ACCENT)

    # 底部截图
    img_path = os.path.join(ASSETS_DIR, "02-job-config-full.png")
    add_image_safe(slide, img_path, Inches(2.5), Inches(6.2), width=Inches(8.0))

    add_page_number(slide, 6)


def slide_07_ai_config(prs):
    """AI 模型配置"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "配置大模型", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.4),
                "可选配置。启用 AI 辅助评估时，需要先完成模型配置和连通性测试", size=14, color=C_GRAY)

    # 左侧截图
    img_path = os.path.join(ASSETS_DIR, "03-api-config-full.png")
    add_image_safe(slide, img_path, Inches(0.5), Inches(1.8), width=Inches(6.5))

    # 右侧配置步骤
    add_textbox(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(0.4),
                "配置 6 步", size=18, bold=True, color=C_DARK)

    steps = [
        "选择服务商（通义千问 / DeepSeek / Kimi / OpenAI ...）",
        "填写 Base URL（系统给出默认地址）",
        "输入 API Key（加密存储在系统钥匙串）",
        "获取或输入模型名称",
        "点击「测试连接」",
        "测试通过 → 保存并添加到列表",
    ]
    for i, step in enumerate(steps):
        y = 2.5 + i * 0.55
        num_dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(y + 0.02), Inches(0.28), Inches(0.28))
        num_dot.fill.solid()
        num_dot.fill.fore_color.rgb = C_BLUE if i < 5 else C_GREEN
        num_dot.line.fill.background()
        tf = num_dot.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        set_font(run, 9, bold=True, color=C_WHITE)

        add_textbox(slide, Inches(7.9), Inches(y - 0.05), Inches(4.8), Inches(0.35),
                    step, size=11, color=C_DGRAY)

    # 底部特性说明
    features = [
        ("🔐 安全存储", "API Key 加密保存在系统钥匙串\n配置文件中不含明文 Key"),
        ("🔄 多接入方式", "同一服务商按 provider+base_url\n区分不同接入方式"),
        ("🔍 模型管理", "搜索/多选/批量测试/新增检测\n支持下线模型提醒"),
    ]

    for i, (title, desc) in enumerate(features):
        x = 0.8 + i * 4.1
        card = add_card(slide, Inches(x), Inches(6.0), Inches(3.8), Inches(1.2))
        add_textbox(slide, Inches(x + 0.2), Inches(6.05), Inches(3.4), Inches(0.3),
                    title, size=12, bold=True, color=C_DARK)
        add_textbox(slide, Inches(x + 0.2), Inches(6.35), Inches(3.4), Inches(0.7),
                    desc, size=10, color=C_GRAY)

    add_page_number(slide, 7)


def slide_08_connection(prs):
    """连接 BOSS 页面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_dark(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
                "连接 BOSS 直聘页面", size=32, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
                "系统通过 Chrome DevTools 协议连接浏览器，操控 BOSS 推荐页面", size=14, color=C_GRAY)

    # 连接步骤
    steps = [
        ("1", "点击「检测/连接浏览器」"),
        ("2", "系统连接或启动 Chrome"),
        ("3", "登录 BOSS 直聘账号"),
        ("4", "打开「推荐牛人」页面"),
        ("5", "回到系统确认状态"),
    ]

    for i, (num, text) in enumerate(steps):
        y = 2.0 + i * 0.7
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.05), Inches(0.4), Inches(0.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_ACCENT
        dot.line.fill.background()
        tf = dot.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, 14, bold=True, color=C_WHITE)
        add_textbox(slide, Inches(1.4), Inches(y), Inches(4.5), Inches(0.4),
                    text, size=16, color=C_WHITE)
        # 连接线
        if i < len(steps) - 1:
            add_rect(slide, Inches(0.97), Inches(y + 0.45), Inches(0.06), Inches(0.3), C_ACCENT)

    # 右侧 — 状态判断表
    add_textbox(slide, Inches(6.5), Inches(2.0), Inches(6), Inches(0.4),
                "浏览器状态判断", size=18, bold=True, color=C_ACCENT)

    statuses = [
        ("未连接", "没有连到 Chrome", "点击检测/连接浏览器", C_RED),
        ("需导航", "已连接但不在推荐页", "手工打开推荐页面", C_YELLOW),
        ("已连接", "已连接到推荐页面", "✅ 可以开始运行", C_GREEN),
    ]

    for i, (status, meaning, action, color) in enumerate(statuses):
        y = 2.8 + i * 1.2
        card = add_rounded_rect(slide, Inches(6.5), Inches(y), Inches(6.0), Inches(1.0),
                                RGBColor(0x24, 0x34, 0x47))
        # 状态标签
        tag = add_rounded_rect(slide, Inches(6.7), Inches(y + 0.15), Inches(1.2), Inches(0.35), color)
        tf = tag.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = status
        set_font(run, 11, bold=True, color=C_WHITE)

        add_textbox(slide, Inches(8.1), Inches(y + 0.1), Inches(4.0), Inches(0.35),
                    meaning, size=12, color=C_WHITE)
        add_textbox(slide, Inches(8.1), Inches(y + 0.5), Inches(4.0), Inches(0.35),
                    "→ " + action, size=11, color=C_GRAY)

    # 底部提示
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.6),
                "⚠️  系统不会替你绕过验证码。出现验证时，先在浏览器完成验证，再回到系统继续。",
                size=13, color=C_YELLOW)

    add_page_number(slide, 8)


def slide_09_running(prs):
    """运行筛选"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "运行筛选", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.4),
                "运行前确认：浏览器已连接 + 推荐页面已打开 + 岗位选择正确", size=14, color=C_GRAY)

    # 左半 — 运行参数 + 策略
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(0.4),
                "运行参数", size=18, bold=True, color=C_DARK)

    params = [
        ("选择岗位", "单岗位优先；全部岗位适合批量"),
        ("滚动轮次", "默认 100，测试可设 20-50"),
        ("AI 辅助评估", "需模型配置，增加耗时和成本"),
        ("自动打招呼", "首次建议\"仅筛选\""),
    ]

    for i, (key, val) in enumerate(params):
        y = 2.4 + i * 0.6
        add_rounded_rect(slide, Inches(0.8), Inches(y), Inches(5.5), Inches(0.5), C_LIGHT)
        add_textbox(slide, Inches(1.0), Inches(y + 0.05), Inches(1.8), Inches(0.35),
                    key, size=12, bold=True, color=C_DARK)
        add_textbox(slide, Inches(2.8), Inches(y + 0.05), Inches(3.3), Inches(0.35),
                    val, size=11, color=C_GRAY)

    # 打招呼策略
    add_textbox(slide, Inches(0.8), Inches(4.9), Inches(5), Inches(0.4),
                "打招呼策略", size=18, bold=True, color=C_DARK)

    strategies = [
        ("仅筛选", "只提取、评分、保存", C_GRAY, "首次测试推荐"),
        ("仅强烈推荐", "≥75 分才打招呼", C_ACCENT, "保守策略"),
        ("推荐+强烈推荐", "≥65 分就打招呼", C_BLUE, "积极策略"),
    ]

    for i, (name, desc, color, note) in enumerate(strategies):
        y = 5.5 + i * 0.6
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y + 0.08), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_textbox(slide, Inches(1.3), Inches(y - 0.02), Inches(2.0), Inches(0.3),
                    name, size=12, bold=True, color=C_DARK)
        add_textbox(slide, Inches(3.2), Inches(y - 0.02), Inches(2.5), Inches(0.3),
                    desc, size=11, color=C_GRAY)
        add_textbox(slide, Inches(5.3), Inches(y - 0.02), Inches(1.5), Inches(0.3),
                    note, size=10, color=color, italic=True)

    # 右侧 — 运行流程图
    img_path = os.path.join(ASSETS_DIR, "04-run-full.png")
    add_image_safe(slide, img_path, Inches(6.8), Inches(1.8), width=Inches(6.0))

    add_page_number(slide, 9)


def slide_10_scoring(prs):
    """评分与推荐等级"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_light(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "评分体系与推荐等级", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.4),
                "四维评分模型 + AI 调整分，最终决定推荐等级", size=14, color=C_GRAY)

    # 四维评分模型卡片
    dims = [
        ("基础分", "25", "固定基础分", C_DARK),
        ("技能匹配", "0~50", "关键词权重 × 命中", C_BLUE),
        ("经验超额", "0~15", "超出要求年限加分", C_ACCENT),
        ("学历档次", "0~10", "硕博/本科/大专分档", RGBColor(0x6C, 0x5C, 0xE7)),
    ]

    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(0.4),
                "四维评分模型（满分 100）", size=16, bold=True, color=C_DARK)

    for i, (name, score, desc, color) in enumerate(dims):
        x = 0.8 + i * 3.1
        card = add_card(slide, Inches(x), Inches(2.4), Inches(2.8), Inches(1.8))
        add_rect(slide, Inches(x), Inches(2.4), Inches(2.8), Inches(0.06), color)
        add_textbox(slide, Inches(x + 0.2), Inches(2.6), Inches(2.4), Inches(0.35),
                    name, size=14, bold=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(3.0), Inches(2.4), Inches(0.5),
                    score, size=28, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(3.6), Inches(2.4), Inches(0.4),
                    desc, size=10, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # 推荐等级
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(5), Inches(0.4),
                "推荐等级", size=16, bold=True, color=C_DARK)

    levels = [
        ("强烈推荐", "≥ 75 分", C_GREEN, "自动打招呼（激进策略）"),
        ("推    荐", "65-74 分", C_BLUE, "自动打招呼（保守策略）"),
        ("待    定", "55-64 分", C_YELLOW, "仅保存，不打招呼"),
        ("淘    汰", "< 55 分", C_RED, "不进入结果统计"),
    ]

    for i, (name, score_range, color, note) in enumerate(levels):
        y = 5.2 + i * 0.55
        # 色块标签
        tag = add_rounded_rect(slide, Inches(0.8), Inches(y), Inches(1.8), Inches(0.4), color)
        tf = tag.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        set_font(run, 12, bold=True, color=C_WHITE)

        add_textbox(slide, Inches(2.8), Inches(y - 0.02), Inches(1.5), Inches(0.35),
                    score_range, size=13, bold=True, color=C_DARK)
        add_textbox(slide, Inches(4.3), Inches(y - 0.02), Inches(4.5), Inches(0.35),
                    note, size=11, color=C_GRAY)

    # AI 调整说明
    card = add_card(slide, Inches(7.5), Inches(4.6), Inches(5.0), Inches(2.6))
    add_textbox(slide, Inches(7.8), Inches(4.7), Inches(4.5), Inches(0.4),
                "🤖 AI 二次评估", size=16, bold=True, color=C_DARK)
    ai_notes = [
        "对 ≥55 分候选人做 LLM 二次评估",
        "调整分范围：±10 分",
        "调整后叠加到规则评分",
        "最多 50 人/次，并发 3 路",
        "含 429 限流退避保护",
    ]
    for i, note in enumerate(ai_notes):
        y = 5.3 + i * 0.4
        add_textbox(slide, Inches(7.8), Inches(y), Inches(4.5), Inches(0.35),
                    f"• {note}", size=11, color=C_DGRAY)

    add_page_number(slide, 10)


def slide_11_results(prs):
    """筛选结果"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "筛选结果", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.4),
                "查看候选人列表、右键操作、单独打招呼、导出 Excel", size=14, color=C_GRAY)

    # 截图
    img_path = os.path.join(ASSETS_DIR, "05-results.png")
    add_image_safe(slide, img_path, Inches(0.5), Inches(1.8), width=Inches(7.5))

    # 右侧操作面板
    add_textbox(slide, Inches(8.5), Inches(1.8), Inches(4), Inches(0.4),
                "常用操作", size=18, bold=True, color=C_DARK)

    ops = [
        ("刷新结果", "从 JSON 重新加载", "🔄"),
        ("导出 Excel", "生成/更新 XLSX", "📊"),
        ("打开 JSON", "用默认程序打开", "📂"),
        ("清空候选人", "操作前自动备份", "🗑"),
    ]

    for i, (name, desc, icon) in enumerate(ops):
        y = 2.5 + i * 0.7
        card = add_card(slide, Inches(8.5), Inches(y), Inches(4.2), Inches(0.55))
        add_textbox(slide, Inches(8.7), Inches(y + 0.05), Inches(0.5), Inches(0.35),
                    icon, size=16, color=C_DARK)
        add_textbox(slide, Inches(9.2), Inches(y + 0.02), Inches(1.8), Inches(0.3),
                    name, size=12, bold=True, color=C_DARK)
        add_textbox(slide, Inches(10.8), Inches(y + 0.05), Inches(1.8), Inches(0.3),
                    desc, size=10, color=C_GRAY)

    # 右键菜单
    add_textbox(slide, Inches(8.5), Inches(5.5), Inches(4), Inches(0.4),
                "右键菜单", size=18, bold=True, color=C_DARK)

    right_click = [
        "查看候选人详情",
        "单独打招呼（定位卡片）",
        "移除候选人",
        "导出选中候选人",
    ]

    for i, item in enumerate(right_click):
        y = 6.1 + i * 0.35
        add_textbox(slide, Inches(8.7), Inches(y), Inches(4.0), Inches(0.3),
                    f"▸ {item}", size=11, color=C_DGRAY)

    add_page_number(slide, 11)


def slide_12_stats(prs):
    """数据统计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_light(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "数据统计", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.4),
                "按岗位和时间范围统计，判断规则质量", size=14, color=C_GRAY)

    # 截图
    img_path = os.path.join(ASSETS_DIR, "06-stats.png")
    add_image_safe(slide, img_path, Inches(0.5), Inches(1.8), width=Inches(6.5))

    # 右侧指标卡片
    metrics = [
        ("总候选人", "≥55 分的人数", C_DARK),
        ("强烈推荐", "≥75 分", C_GREEN),
        ("推荐", "65-74 分", C_BLUE),
        ("待定", "55-64 分", C_YELLOW),
        ("已打招呼", "已发送消息数", C_ACCENT),
        ("优质率", "强推+推荐占比", RGBColor(0x6C, 0x5C, 0xE7)),
        ("打招呼率", "已打招呼占比", RGBColor(0x00, 0xB8, 0x94)),
        ("平均分", "当前范围均分", C_GRAY),
    ]

    add_textbox(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(0.4),
                "统计指标", size=18, bold=True, color=C_DARK)

    for i, (name, desc, color) in enumerate(metrics):
        row = i // 2
        col = i % 2
        x = 7.5 + col * 2.7
        y = 2.4 + row * 1.15
        card = add_card(slide, Inches(x), Inches(y), Inches(2.5), Inches(1.0))
        add_rect(slide, Inches(x), Inches(y), Inches(0.06), Inches(1.0), color)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.1), Inches(2.2), Inches(0.3),
                    name, size=13, bold=True, color=C_DARK)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.5), Inches(2.2), Inches(0.35),
                    desc, size=10, color=C_GRAY)

    # 底部洞察
    add_rounded_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), C_WHITE)
    add_textbox(slide, Inches(1.2), Inches(6.85), Inches(11.0), Inches(0.4),
                "💡 优质率极低？通常不是市场没人，而是规则过硬、关键词偏窄或推荐池不匹配。",
                size=12, color=C_DGRAY)

    add_page_number(slide, 12)


def slide_13_best_practices(prs):
    """推荐操作习惯"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "推荐操作习惯", size=32, bold=True, color=C_DARK)
    add_textbox(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.4),
                "新岗位先试跑，成熟岗位日常跑", size=14, color=C_GRAY)

    # 左半 — 新岗位
    card1 = add_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.0))
    add_rect(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(0.06), C_ACCENT)
    tag1 = add_rounded_rect(slide, Inches(0.9), Inches(2.1), Inches(2.5), Inches(0.4), C_ACCENT)
    tf = tag1.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "新岗位首次跑"
    set_font(run, 13, bold=True, color=C_WHITE)

    new_steps = [
        "先只配置规则，不启用 AI",
        "打招呼策略选「不打招呼（仅筛选）」",
        "滚动轮次设为 20-50（少量测试）",
        "看结果里的误杀、误放",
        "调整岗位规则",
        "确认规则后开启「仅强烈推荐」",
    ]
    for i, step in enumerate(new_steps):
        y = 2.8 + i * 0.6
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.03), Inches(0.32), Inches(0.32))
        num.fill.solid()
        num.fill.fore_color.rgb = C_ACCENT
        num.line.fill.background()
        ntf = num.text_frame
        p = ntf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        nr = p.add_run()
        nr.text = str(i + 1)
        set_font(nr, 11, bold=True, color=C_WHITE)
        add_textbox(slide, Inches(1.5), Inches(y), Inches(4.5), Inches(0.35),
                    step, size=12, color=C_DGRAY)

    # 右半 — 成熟岗位
    card2 = add_card(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0))
    add_rect(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(0.06), C_BLUE)
    tag2 = add_rounded_rect(slide, Inches(7.2), Inches(2.1), Inches(2.5), Inches(0.4), C_BLUE)
    tf2 = tag2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "成熟岗位日常"
    set_font(run2, 13, bold=True, color=C_WHITE)

    daily_steps = [
        ("检查浏览器连接", "✅"),
        ("确认 BOSS 当前职位", "🎯"),
        ("运行筛选", "▶"),
        ("查看筛选结果", "📋"),
        ("导出 Excel", "📊"),
        ("查看统计，判断是否微调规则", "🔍"),
    ]
    for i, (step, icon) in enumerate(daily_steps):
        y = 2.8 + i * 0.6
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), Inches(y + 0.03), Inches(0.32), Inches(0.32))
        num.fill.solid()
        num.fill.fore_color.rgb = C_BLUE
        num.line.fill.background()
        ntf = num.text_frame
        p = ntf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        nr = p.add_run()
        nr.text = str(i + 1)
        set_font(nr, 11, bold=True, color=C_WHITE)
        add_textbox(slide, Inches(7.7), Inches(y - 0.02), Inches(0.4), Inches(0.35),
                    icon, size=14, color=C_DARK)
        add_textbox(slide, Inches(8.2), Inches(y), Inches(4.0), Inches(0.35),
                    step, size=12, color=C_DGRAY)

    add_page_number(slide, 13)


def slide_14_faq(prs):
    """常见问题速查"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_light(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.6), height=Inches(0.5))
    add_textbox(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.6),
                "常见问题速查", size=32, bold=True, color=C_DARK)

    faqs = [
        ("提示「未连接」",
         "回到运行控制 → 点击检测/连接浏览器 → 确认 Chrome 已启动且打开推荐页",
         C_RED),
        ("提示「需导航」",
         "Chrome 已连接但不在推荐页 → 手工打开目标岗位的推荐牛人页面",
         C_YELLOW),
        ("BOSS 弹出验证码",
         "系统暂停等待 → 在浏览器完成验证 → 回到系统选择继续",
         C_YELLOW),
        ("AI 评估失败",
         "检查 API Key → Base URL → 模型名称 → 是否已开通 → 额度是否耗尽（不要先改规则）",
         C_BLUE),
        ("候选人很少/质量差",
         "检查：职位选对？地点过窄？必要过硬？关键词过少？轮次太低？推荐池枯竭？",
         C_ACCENT),
        ("Excel 没有更新",
         "筛选结果 → 刷新结果 → 导出 Excel → 确认 JSON 有新数据",
         C_GRAY),
    ]

    for i, (question, answer, color) in enumerate(faqs):
        y = 1.5 + i * 0.95
        card = add_card(slide, Inches(0.8), Inches(y), Inches(11.7), Inches(0.8))
        # 左侧色条
        add_rect(slide, Inches(0.8), Inches(y), Inches(0.06), Inches(0.8), color)
        # 问题编号
        q_num = add_rounded_rect(slide, Inches(1.1), Inches(y + 0.15), Inches(0.5), Inches(0.5), color)
        tf = q_num.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Q" + str(i + 1)
        set_font(run, 11, bold=True, color=C_WHITE)

        add_textbox(slide, Inches(1.8), Inches(y + 0.05), Inches(3.0), Inches(0.35),
                    question, size=13, bold=True, color=C_DARK)
        add_textbox(slide, Inches(1.8), Inches(y + 0.4), Inches(10.3), Inches(0.35),
                    answer, size=11, color=C_GRAY)

    add_page_number(slide, 14)


def slide_15_checklist(prs):
    """快速检查清单 + 封底"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_dark(slide)

    # 装饰
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x24, 0x34, 0x47)
    circle.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
                "快速检查清单", size=32, bold=True, color=C_WHITE)

    # 运行前检查
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
                "运行前", size=18, bold=True, color=C_ACCENT)

    pre_checks = [
        "岗位规则已保存",
        "BOSS 网页端已登录",
        "Chrome 已连接",
        "当前页面是目标岗位推荐牛人页面",
        "系统选择岗位与 BOSS 页面一致",
        "AI 模型已测试通过",
        "新岗位未直接启用大范围打招呼",
    ]

    for i, item in enumerate(pre_checks):
        y = 2.1 + i * 0.45
        add_textbox(slide, Inches(1.0), Inches(y), Inches(5.5), Inches(0.35),
                    f"☑  {item}", size=12, color=C_WHITE)

    # 运行后检查
    add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5), Inches(0.4),
                "运行后", size=18, bold=True, color=C_ACCENT)

    post_checks = [
        "筛选结果页已刷新",
        "Excel 已导出",
        "数据统计已查看",
        "异常候选人已核对",
        "必要时备份 JSON",
    ]

    for i, item in enumerate(post_checks):
        y = 2.1 + i * 0.45
        add_textbox(slide, Inches(7.2), Inches(y), Inches(5.5), Inches(0.35),
                    f"☑  {item}", size=12, color=C_WHITE)

    # 底部数据文件
    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(12), Inches(0.4),
                "关键数据文件", size=16, bold=True, color=C_GRAY)

    files = [
        ("job_config.json", "岗位规则"),
        ("api_config.json", "模型配置"),
        ("candidates_all.json", "候选人数据"),
        ("candidates_all.xlsx", "Excel 导出"),
        ("selectors.json", "页面选择器"),
    ]

    for i, (fname, fdesc) in enumerate(files):
        x = 0.8 + i * 2.5
        add_rounded_rect(slide, Inches(x), Inches(6.2), Inches(2.3), Inches(0.6),
                         RGBColor(0x24, 0x34, 0x47))
        add_textbox(slide, Inches(x + 0.1), Inches(6.2), Inches(2.1), Inches(0.3),
                    fname, size=10, bold=True, color=C_ACCENT)
        add_textbox(slide, Inches(x + 0.1), Inches(6.5), Inches(2.1), Inches(0.25),
                    fdesc, size=9, color=C_GRAY)

    # 最底部
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(12), Inches(0.4),
                "BOSS 招聘系统 v2.9 · 操作指南 · 2026",
                size=11, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 15)


# ─── 主流程 ──────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        slide_01_cover,
        slide_02_overview,
        slide_03_workflow,
        slide_04_navigation,
        slide_05_job_config,
        slide_06_jd_parse,
        slide_07_ai_config,
        slide_08_connection,
        slide_09_running,
        slide_10_scoring,
        slide_11_results,
        slide_12_stats,
        slide_13_best_practices,
        slide_14_faq,
        slide_15_checklist,
    ]

    for fn in slides:
        fn(prs)
        print(f"  [OK] {fn.__name__}")

    prs.save(OUTPUT_PATH)
    print(f"\nDone: {OUTPUT_PATH}")
    print(f"  Total: {len(slides)} slides")


if __name__ == "__main__":
    main()
